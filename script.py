#!/usr/bin/env python3
import datetime
import csv
import sys
import random
import os
import subprocess
import re
import shutil
import json
import zipfile
import webbrowser  # Aggiunto per aprire la documentazione
from pathlib import Path
from collections import defaultdict, Counter

# --- Gestione codifica su Windows ---
if sys.platform == 'win32':
    try:
        subprocess.run(['chcp', '65001'], check=True, capture_output=True, text=True,
                       creationflags=subprocess.CREATE_NO_WINDOW)
    except (FileNotFoundError, subprocess.CalledProcessError):
        pass

# --- IMPOSTAZIONE SEPARATORE ---
FILENAME_SEPARATOR = "--"

# --- GESTIONE DIPENDENZE OPZIONALI ---
try:
    if not sys.stdout.isatty():
        raise ImportError("Non è un terminale interattivo, fallback a input standard.")
    import inquirer
except ImportError:
    inquirer = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    Presentation = None

# Aggiunto import per Pillow per una gestione robusta delle immagini
try:
    from PIL import Image
except ImportError:
    print(
        "AVVISO: La libreria Pillow non è installata (pip install Pillow). Le funzioni PPTX potrebbero non funzionare.")
    Image = None


# --- FUNZIONI DI UTILITÀ E INPUT ---

def _sanitize_name(name: str) -> str:
    invalid = '<>:"/\\|?*'
    sanitized = ''.join('_' if c in invalid else c for c in name)
    sanitized = sanitized.strip().rstrip(' .')
    while '  ' in sanitized:
        sanitized = sanitized.replace('  ', ' ')
    return sanitized or 'Senza_Nome'


def _exit_requested(s: str) -> bool:
    return s.strip().lower() in {"q", "quit", "esci", "exit", "stop"}


def _input_or_exit(prompt: str) -> str:
    try:
        val = input(prompt)
    except (EOFError, KeyboardInterrupt):
        print("\nUscita richiesta.")
        sys.exit(0)
    if _exit_requested(val):
        print("Uscita richiesta.")
        sys.exit(0)
    return val


def _confirm_or_exit(message: str, default: bool = False) -> bool:
    if inquirer:
        full_message = f"{message.strip()} (premi Ctrl+C per uscire)"
        questions = [inquirer.Confirm('confirm', message=full_message, default=default)]
        try:
            answers = inquirer.prompt(questions)
            if answers is None: sys.exit(0)
            return answers['confirm']
        except (KeyboardInterrupt, EOFError):
            print("\nUscita richiesta.");
            sys.exit(0)
    else:
        prompt = f"{message.strip()} [s/n]: "
        resp = _input_or_exit(prompt).strip().lower()
        if resp == "": return default
        return resp in {"s", "si", "sì", "y", "yes"}


def _open_folder(path: Path) -> None:
    try:
        if os.name == 'nt':
            os.startfile(str(path))
        elif sys.platform == 'darwin':
            subprocess.run(['open', str(path)])
        else:
            subprocess.run(['xdg-open', str(path)])
    except Exception as e:
        print(f"Errore durante l'apertura della cartella '{path}': {e}")


def _make_snake(s: str) -> str:
    snake = "".join(ch.lower() if ch.isalnum() else "_" for ch in s)
    while "__" in snake: snake = snake.replace("__", "_")
    return snake.strip("_")


def _parse_date_flexible(date_str: str) -> datetime.date | None:
    if not date_str:
        return None
    try:
        return datetime.datetime.strptime(date_str, '%d/%m/%Y').date()
    except ValueError:
        try:
            return datetime.datetime.strptime(date_str, '%Y-%m-%d').date()
        except ValueError:
            return None


# --- GESTIONE LOG CONCORSI ---

def _get_log_path(root_dir: Path) -> Path:
    return root_dir / "ContestRegistry.json"


def _read_log(log_path: Path) -> dict:
    if not log_path.exists():
        return {"contests": {}}
    try:
        with log_path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, FileNotFoundError, PermissionError) as e:
        print(
            f"⚠️  ATTENZIONE: Impossibile leggere il file di log '{log_path}'. Verrà usato un log vuoto. Dettagli: {e}")
        return {"contests": {}}


def _write_log(log_path: Path, data: dict) -> None:
    try:
        with log_path.open("w", encoding="utf-8") as f:
            json.dump(data, f, indent=2)
    except (IOError, PermissionError) as e:
        print("\n" + "---" * 15)
        print(f"❌ ERRORE CRITICO: Impossibile scrivere nel file di log '{log_path}'.")
        print(f"   Dettagli tecnici: {e}")
        print("\n   Questo errore (PermissionError) solitamente accade per uno di questi motivi:")
        print("   1. Permessi Mancanti: Non hai i permessi necessari per scrivere in questa cartella.")
        print(
            "   2. File Bloccato: Un altro programma (come un editor di testo o un antivirus) sta tenendo il file aperto.")
        _input_or_exit("\nPremi Invio per chiudere il programma.")
        sys.exit(1)


def _update_contest_in_log(log_path: Path, contest_name: str, deadline: str | None):
    log_data = _read_log(log_path)
    if contest_name not in log_data["contests"]:
        log_data["contests"][contest_name] = {}
    log_data["contests"][contest_name]["deadline"] = deadline
    _write_log(log_path, log_data)


# --- FUNZIONI DI LETTURA/SCRITTURA FILE ---

def _read_existing_criteria(csv_path: Path) -> list[str]:
    if not csv_path.exists(): return []
    try:
        with csv_path.open("r", encoding="utf-8-sig", newline="") as f:
            reader = csv.reader(f, delimiter=';')
            return [c for c in next(reader, [])[1:] if c]
    except Exception:
        return []


def _read_existing_titles(csv_path: Path) -> list[str]:
    if not csv_path.exists(): return []
    titles = []
    try:
        with csv_path.open("r", encoding="utf-8-sig", newline="") as f:
            reader = csv.reader(f, delimiter=';')
            next(reader, None)
            for row in reader:
                if row: titles.append(row[0])
    except Exception:
        pass
    return titles


def _write_consolidated_csv(csv_path: Path, criteria: list[str], titles: list[str]) -> None:
    header = [""] + criteria
    rows = [header] + [[t] + [""] * len(criteria) for t in titles]
    try:
        csv_path.parent.mkdir(parents=True, exist_ok=True)
        with csv_path.open("w", newline="", encoding="utf-8-sig") as f:
            writer = csv.writer(f, delimiter=';')
            writer.writerows(rows)
    except (IOError, PermissionError) as e:
        print(f"❌ ERRORE: Impossibile scrivere il file CSV '{csv_path.name}'. Dettagli: {e}")


def _normalize_string(text: str) -> str:
    text = text.replace('_', ' ').replace('-', ' ')
    return ' '.join(text.split()).strip()


def _parse_filename(path: Path) -> tuple[str, str]:
    stem = path.stem.strip()
    raw_author, raw_title = "", ""
    if FILENAME_SEPARATOR in stem:
        parts = stem.split(FILENAME_SEPARATOR, 1)
        raw_author, raw_title = parts[0], parts[1]
    else:
        match = re.search(r'_(?P<num>\d+)_', stem)
        if match:
            raw_author, raw_title = stem[:match.start()], stem[match.end():]
        else:
            raw_title = stem

    author = _normalize_string(raw_author).title()
    normalized_title = _normalize_string(raw_title)
    title = normalized_title[0].upper() + normalized_title[1:] if normalized_title else ""
    return author or "Autore Sconosciuto", title or "Titolo Sconosciuto"


# --- DASHBOARD DI STATO ---

def _display_contest_status(base: Path, folder_name: str):
    pictures_path = base / "pictures"
    judges_dir = base / "judges"
    csv_path = base / "jury_kit" / f"{_make_snake(folder_name)} giuria.csv"
    log_path = _get_log_path(base.parent)
    log_data = _read_log(log_path)
    contest_info = log_data.get("contests", {}).get(folder_name, {})
    deadline_str = contest_info.get("deadline")
    deadline_date = _parse_date_flexible(deadline_str)
    photo_count = len([p for p in pictures_path.iterdir() if p.is_file() and p.suffix.lower() in {'.jpg', '.jpeg'}])
    criteria = _read_existing_criteria(csv_path)
    jurors = [d.name for d in judges_dir.iterdir() if d.is_dir()]
    submitted_votes = [p.parent.name for p in _find_judge_csvs(judges_dir)]

    print("\n" + "---" * 15)
    print(f"📊 STATO CONCORSO: {folder_name}")
    if deadline_date:
        is_expired = deadline_date < datetime.date.today()
        print(f"  - Scadenza:           {deadline_str} {'(Scaduto)' if is_expired else ''}")
    else:
        print("  - Scadenza:           Non impostata")
    print("---" * 15)
    print(f"  - Foto Trovate:       {photo_count}")
    print(f"  - Criteri Definiti:   {len(criteria)}")
    print(f"  - Giurati Registrati: {len(jurors)}")
    print(f"  - Voti Consegnati:    {len(submitted_votes)} su {len(jurors)}")
    print("---" * 15)


# --- FUNZIONI AZIONE DEL MENU ---

def _action_manage_criteria(base: Path, folder_name: str):
    """Gestisce i criteri di valutazione con un menu interattivo (Aggiungi/Rinomina/Elimina)."""
    csv_path = base / "jury_kit" / f"{_make_snake(folder_name)} giuria.csv"

    if not csv_path.exists() or not _read_existing_criteria(csv_path):
        print("\nNessun criterio definito per questo concorso.")
        if _confirm_or_exit("Vuoi iniziare con una lista di criteri standard?", default=True):
            default_criteria = ["Attinenza al tema", "Tecnica", "Creatività", "Impressione generale"]
            _write_consolidated_csv(csv_path, default_criteria, [])
            print("✅ Criteri standard aggiunti.")
        else:
            _write_consolidated_csv(csv_path, [], [])

    while True:
        print("\n### 📝 Gestione Criteri di Valutazione ###")
        criteria = _read_existing_criteria(csv_path)

        if criteria:
            print("\nCriteri attuali:")
            for criterion in criteria:
                print(f"  - {criterion}")
        else:
            print("\nAl momento non ci sono criteri definiti.")

        choices = ["Aggiungi un nuovo criterio"]
        if criteria:
            choices.extend(["Rinomina un criterio", "Elimina un criterio"])
        choices.append("Torna al menu principale")

        action = ""
        print("\n--- Azioni Disponibili ---")
        if inquirer:
            questions = [inquirer.List('action', message="Scegli un'azione", choices=choices, carousel=True)]
            answers = inquirer.prompt(questions)
            if not answers: return
            action = answers['action']
        else:
            for i, choice in enumerate(choices, 1): print(f"  [{i}] {choice}")
            while True:
                try:
                    choice_num = int(_input_or_exit("Scelta: "))
                    if 1 <= choice_num <= len(choices):
                        action = choices[choice_num - 1]
                        break
                    else:
                        print("Scelta non valida.")
                except ValueError:
                    print("Inserisci un numero.")

        titles = _read_existing_titles(csv_path)

        if "Aggiungi" in action:
            new_criterion = _input_or_exit("Nome del nuovo criterio (vuoto per terminare): ").strip()
            if new_criterion:
                criteria.append(new_criterion)
                _write_consolidated_csv(csv_path, criteria, titles)
                print(f"✅ Criterio '{new_criterion}' aggiunto.")

        elif "Rinomina" in action and criteria:
            old_criterion, choice_idx = "", -1
            print("\nCriteri esistenti:")
            for i, c in enumerate(criteria, 1): print(f"  [{i}] {c}")

            if inquirer:
                q = [inquirer.List('old', message="Quale criterio vuoi rinominare?", choices=criteria)]
                answers = inquirer.prompt(q)
                if answers: old_criterion = answers['old']
            else:
                while True:
                    try:
                        choice_idx = int(_input_or_exit("Numero del criterio da rinominare: ")) - 1
                        if 0 <= choice_idx < len(criteria):
                            old_criterion = criteria[choice_idx]
                            break
                        else:
                            print("Numero non valido.")
                    except (ValueError, IndexError):
                        print("Inserisci un numero valido.")

            if old_criterion:
                new_criterion = _input_or_exit(f"Nuovo nome per '{old_criterion}': ").strip()
                if new_criterion:
                    idx_to_rename = criteria.index(old_criterion)
                    criteria[idx_to_rename] = new_criterion
                    _write_consolidated_csv(csv_path, criteria, titles)
                    print(f"✅ Criterio rinominato in '{new_criterion}'.")

        elif "Elimina" in action and criteria:
            to_delete = ""
            print("\nCriteri esistenti:")
            for i, c in enumerate(criteria, 1): print(f"  [{i}] {c}")

            if inquirer:
                q = [inquirer.List('del', message="Quale criterio vuoi eliminare?", choices=criteria)]
                answers = inquirer.prompt(q)
                if answers: to_delete = answers['del']
            else:
                while True:
                    try:
                        choice_idx = int(_input_or_exit("Numero del criterio da eliminare: ")) - 1
                        if 0 <= choice_idx < len(criteria):
                            to_delete = criteria[choice_idx]
                            break
                        else:
                            print("Numero non valido.")
                    except (ValueError, IndexError):
                        print("Inserisci un numero valido.")

            if to_delete and _confirm_or_exit(f"Sei sicuro di voler eliminare il criterio '{to_delete}'?",
                                              default=False):
                criteria.remove(to_delete)
                _write_consolidated_csv(csv_path, criteria, titles)
                print(f"🗑️ Criterio '{to_delete}' eliminato.")

        elif "Torna" in action:
            break


def _action_generate_kit(base: Path, folder_name: str):
    print("\n### 🔄 Genera/Aggiorna Jury Kit Completo ###")
    pictures_path = base / "pictures"
    jury_kit_dir = base / "jury_kit"
    csv_path = jury_kit_dir / f"{_make_snake(folder_name)} giuria.csv"

    criteria = _read_existing_criteria(csv_path)
    if not criteria:
        print("\n❌ ERRORE: Nessun criterio di valutazione definito.")
        print("   Usa l'opzione 'Gestisci Criteri di Valutazione' prima di generare il kit.")
        _input_or_exit("Premi Invio per continuare...")
        return

    photo_paths = [p for p in pictures_path.iterdir() if p.is_file() and p.suffix.lower() in {'.jpg', '.jpeg'}]
    if not photo_paths:
        print("\n⚠️ ATTENZIONE: Nessuna foto trovata nella cartella 'pictures'.")
        if not _confirm_or_exit("Vuoi continuare a generare un kit senza foto?", default=False):
            print("Operazione annullata.")
            return

    # <<< CONTROLLO DUPLICATI >>>
    parsed_titles = [_parse_filename(p)[1] for p in photo_paths]
    title_counts = Counter(parsed_titles)
    duplicates = [title for title, count in title_counts.items() if count > 1]

    if duplicates:
        print("\n" + "---" * 15)
        print("❌ ERRORE CRITICO: Trovati titoli di foto duplicati!")
        print("   La generazione del kit è bloccata perché più file producono lo stesso titolo dopo il parsing.")
        print("\n   Titoli duplicati rilevati:")
        for title in duplicates:
            print(f"     - \"{title}\"")
        print("\n   SOLUZIONE: Rinomina i file originali nella cartella 'pictures' per garantire che ogni foto")
        print("   abbia un titolo unico e riprova.")
        print("---" * 15)
        _input_or_exit("\nPremi Invio per tornare al menu...")
        return

    current_titles = sorted(parsed_titles)  # Riusiamo la lista già calcolata
    old_titles = _read_existing_titles(csv_path)

    if sorted(current_titles) == sorted(old_titles):
        print("\nNessuna nuova foto rilevata.")
        if not _confirm_or_exit("Vuoi rigenerare comunque tutti i file del kit (PPT, ZIP, etc.)?", default=False):
            print("Operazione annullata.")
            return

    print("\n--- Riepilogo Azioni ---")
    print(f"  - Criteri da usare: {', '.join(criteria)}")
    print(f"  - Foto da includere: {len(current_titles)}")
    print("  - Verranno generati/aggiornati: CSV (master), PPT, cartella foto e ZIP (senza CSV).")
    print("--------------------------")

    if not _confirm_or_exit("Procedere con la generazione del kit?", default=True):
        print("❌ Operazione annullata.")
        return

    titles_for_csv = current_titles
    random.shuffle(titles_for_csv)

    print("\nAggiornamento in corso...")
    _write_consolidated_csv(csv_path, criteria, titles_for_csv)
    print(f"✅ File master CSV '{csv_path.name}' aggiornato.")

    if photo_paths:
        ppt_path = jury_kit_dir / f"{_make_snake(folder_name)}.pptx"
        parsed_photos = {_parse_filename(p)[1]: p for p in photo_paths}
        ordered_pairs = [(t, parsed_photos[t]) for t in titles_for_csv if t in parsed_photos]

        _build_ppt(ppt_path, ordered_pairs)
        print(f"✅ Presentazione '{ppt_path.name}' aggiornata.")

        _update_original_pictures_folder(jury_kit_dir, ordered_pairs)
        print(f"✅ Cartella 'original_pictures' aggiornata.")

    archive_base_path = base / f"{_make_snake(folder_name)}_jury_kit"
    _create_jury_kit_zip(jury_kit_dir, archive_base_path)
    print("✅ Kit per la giuria generato con successo!")

    _input_or_exit("Premi Invio per tornare al menu...")


def _action_manage_jurors(judges_dir: Path):
    while True:
        print("\n### 🛠️  Gestione Giurati ###")
        jurors = sorted([d.name for d in judges_dir.iterdir() if d.is_dir()])

        if jurors:
            print("\nGiurati registrati:")
            for juror in jurors:
                print(f"  - {juror}")
        else:
            print("\nNessun giurato registrato.")

        choices = ["Aggiungi un nuovo giurato"]
        if jurors:
            choices.extend(["Rinomina un giurato esistente", "Elimina un giurato"])
        choices.append("Torna al menu principale")

        action = ""
        print("\n--- Azioni Disponibili ---")
        if inquirer:
            questions = [inquirer.List('action', message="Scegli un'azione", choices=choices, carousel=True)]
            answers = inquirer.prompt(questions)
            if not answers: return
            action = answers['action']
        else:
            for i, choice in enumerate(choices, 1): print(f"  [{i}] {choice}")
            while True:
                try:
                    choice_num = int(_input_or_exit("Scelta: "))
                    if 1 <= choice_num <= len(choices):
                        action = choices[choice_num - 1]
                        break
                    else:
                        print("Scelta non valida.")
                except ValueError:
                    print("Inserisci un numero.")

        if "Aggiungi" in action:
            while True:
                name = _input_or_exit("Nome nuovo giurato (vuoto per terminare): ").strip()
                if not name: break
                safe_name = _sanitize_name(name)
                if (judges_dir / safe_name).exists():
                    print(f"ATTENZIONE: Giurato '{safe_name}' esiste già.")
                else:
                    (judges_dir / safe_name).mkdir(parents=True, exist_ok=True)
                    print(f"✅ Creata cartella per: {safe_name}")
        elif "Rinomina" in action and jurors:
            old_name = ""
            print("\nGiurati esistenti:")
            for i, name in enumerate(jurors, 1): print(f"  [{i}] {name}")

            if inquirer:
                q = [inquirer.List('old_name', message="Quale giurato vuoi rinominare?", choices=jurors)]
                answers = inquirer.prompt(q)
                if answers: old_name = answers['old_name']
            else:
                while True:
                    try:
                        idx = int(_input_or_exit("Numero del giurato da rinominare: ")) - 1
                        if 0 <= idx < len(jurors):
                            old_name = jurors[idx]
                            break
                        else:
                            print("Scelta non valida.")
                    except (ValueError, IndexError):
                        print("Inserisci un numero valido.")

            if old_name:
                new_name_raw = _input_or_exit(f"Nuovo nome per '{old_name}': ").strip()
                if new_name_raw:
                    new_name = _sanitize_name(new_name_raw)
                    if (judges_dir / new_name).exists():
                        print(f"ERRORE: Esiste già un giurato di nome '{new_name}'.")
                    else:
                        try:
                            (judges_dir / old_name).rename(judges_dir / new_name)
                            print(f"✅ '{old_name}' rinominato in '{new_name}'.")
                        except OSError as e:
                            print(f"❌ ERRORE durante la rinomina: {e}")
        elif "Elimina" in action and jurors:
            name_to_delete = ""
            print("\nGiurati esistenti:")
            for i, name in enumerate(jurors, 1): print(f"  [{i}] {name}")

            if inquirer:
                q = [inquirer.List('del_name', message="Quale giurato vuoi eliminare?", choices=jurors)]
                answers = inquirer.prompt(q)
                if answers: name_to_delete = answers['del_name']
            else:
                while True:
                    try:
                        idx = int(_input_or_exit("Numero del giurato da eliminare: ")) - 1
                        if 0 <= idx < len(jurors):
                            name_to_delete = jurors[idx]
                            break
                        else:
                            print("Scelta non valida.")
                    except (ValueError, IndexError):
                        print("Inserisci un numero valido.")

            if name_to_delete and _confirm_or_exit(f"Sei sicuro di voler eliminare '{name_to_delete}' e i suoi voti?",
                                                   default=False):
                try:
                    shutil.rmtree(judges_dir / name_to_delete)
                    print(f"🗑️  Giurato '{name_to_delete}' eliminato.")
                except OSError as e:
                    print(f"❌ ERRORE durante l'eliminazione: {e}")
        elif "Torna" in action:
            break


def _action_generate_leaderboard(base: Path, folder_name: str):
    print("\n### 🏆 Generazione Classifica Finale ###")
    if Presentation is None:
        print("\n❌ ERRORE: La libreria 'python-pptx' è necessaria. Installala con: pip install python-pptx")
        _input_or_exit("Premi Invio per tornare al menu...")
        return

    csv_path = base / "jury_kit" / f"{_make_snake(folder_name)} giuria.csv"
    criteria = _read_existing_criteria(csv_path)
    juror_dirs = [d for d in (base / "judges").iterdir() if d.is_dir()]
    submitted_csvs = _find_judge_csvs(base / "judges")

    if not all([criteria, juror_dirs, submitted_csvs]):
        print("Impossibile generare la classifica. Condizioni non soddisfatte:")
        if not criteria: print("  - ❌ Nessun criterio definito.")
        if not juror_dirs: print("  - ❌ Nessun giurato registrato.")
        if not submitted_csvs: print("  - ❌ Nessun giurato ha consegnato i voti.")
        _input_or_exit("\nPremi Invio per tornare al menu...")
        return

    _generate_leaderboard_logic(base, csv_path)
    _open_folder(base / "leaderboard")
    _input_or_exit("Premi Invio per tornare al menu...")


def _action_edit_deadline(root_dir: Path, contest_name: str):
    print("\n### ✏️  Modifica Data di Scadenza ###")
    log_path = _get_log_path(root_dir)
    log_data = _read_log(log_path)
    current_deadline = log_data.get("contests", {}).get(contest_name, {}).get("deadline")
    print(f"Data di scadenza attuale: {current_deadline or 'Non impostata'}")

    while True:
        new_date_str = _input_or_exit("Nuova data (GG/MM/AAAA) o vuoto per rimuovere: ").strip()
        if not new_date_str:
            _update_contest_in_log(log_path, contest_name, None)
            print("✅ Data di scadenza rimossa.")
            break
        try:
            datetime.datetime.strptime(new_date_str, '%d/%m/%Y')
            _update_contest_in_log(log_path, contest_name, new_date_str)
            print(f"✅ Data di scadenza aggiornata a {new_date_str}.")
            break
        except ValueError:
            print("Formato data non valido. Usa GG/MM/AAAA. Riprova.")
    _input_or_exit("Premi Invio per tornare al menu...")


# --- LOGICA PRINCIPALE (CORE) ---

def _select_or_create_contest(root_dir: Path) -> tuple[Path, str] | None:
    """Seleziona un concorso, ne crea uno nuovo o permette di uscire."""
    log_path = _get_log_path(root_dir)
    log_data = _read_log(log_path)
    today = datetime.date.today()
    existing_contests = sorted([d.name for d in root_dir.iterdir() if d.is_dir()])
    active_contests, expired_contests = [], []

    for contest_name in existing_contests:
        info = log_data.get("contests", {}).get(contest_name, {})
        deadline_str = info.get("deadline")
        deadline_date = _parse_date_flexible(deadline_str)
        if deadline_date and deadline_date < today:
            expired_contests.append(contest_name)
        else:
            active_contests.append(contest_name)

    new_opt = ">> CREA NUOVO CONCORSO <<"
    exit_opt = "🚪 Esci dal programma"
    active_sep = "--- Concorsi Attivi ---"
    expired_sep = "--- Concorsi Scaduti ---"
    separators = [active_sep, expired_sep]

    choices = [new_opt]
    if active_contests: choices.extend([active_sep] + active_contests)
    if expired_contests: choices.extend([expired_sep] + expired_contests)
    choices.append(exit_opt)

    chosen_name = ""
    if inquirer:
        q = [inquirer.List('contest', message="Scegli un'azione o un concorso", choices=choices, carousel=True)]
        while not chosen_name:
            answers = inquirer.prompt(q)
            if not answers: return None
            if answers['contest'] not in separators:
                chosen_name = answers['contest']
    else:
        print("\nScegli un'opzione:")
        display_map = {}
        current_idx = 1

        print(f"  [{current_idx}] {new_opt}")
        display_map[current_idx] = new_opt
        current_idx += 1

        if active_contests:
            print("\n--- Concorsi Attivi ---")
            for name in active_contests:
                print(f"  [{current_idx}] {name}")
                display_map[current_idx] = name
                current_idx += 1

        if expired_contests:
            print("\n--- Concorsi Scaduti ---")
            for name in expired_contests:
                print(f"  [{current_idx}] {name}")
                display_map[current_idx] = name
                current_idx += 1

        print(f"\n  [0] {exit_opt}")

        while True:
            try:
                choice = int(_input_or_exit("\nScelta: "))
                if choice == 0:
                    chosen_name = exit_opt
                    break
                elif choice in display_map:
                    chosen_name = display_map[choice]
                    break
                else:
                    print("Scelta non valida.")
            except ValueError:
                print("Inserisci un numero.")

    if chosen_name == exit_opt:
        return None

    if chosen_name == new_opt:
        name = _input_or_exit("Nome del nuovo concorso: ")
        safe_name = _sanitize_name(name)
        deadline_str, deadline = "", None
        while True:
            deadline_str = _input_or_exit("Data di scadenza (GG/MM/AAAA, opzionale): ").strip()
            if not deadline_str: break
            try:
                datetime.datetime.strptime(deadline_str, '%d/%m/%Y')
                deadline = deadline_str
                break
            except ValueError:
                print("Formato data non valido. Riprova.")
        _update_contest_in_log(log_path, safe_name, deadline)
        return root_dir / safe_name, safe_name
    else:
        return root_dir / chosen_name, chosen_name


def _update_original_pictures_folder(jury_kit_dir: Path, title_path_pairs: list[tuple[str, Path]]):
    original_pictures_dir = jury_kit_dir / "original_pictures"
    try:
        if original_pictures_dir.exists(): shutil.rmtree(original_pictures_dir)
        original_pictures_dir.mkdir(exist_ok=True)
        for title, original_path in title_path_pairs:
            safe_title = _sanitize_name(title)
            new_filename = f"{safe_title}{original_path.suffix}"
            shutil.copy(original_path, original_pictures_dir / new_filename)
    except Exception as e:
        print(f"ERRORE durante l'aggiornamento di 'original_pictures': {e}")


def _create_jury_kit_zip(jury_kit_dir: Path, archive_base_name: Path):
    """Crea un archivio ZIP della cartella jury_kit, escludendo i file CSV."""
    zip_path = archive_base_name.with_suffix('.zip')
    try:
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for item in jury_kit_dir.rglob('*'):
                if item.suffix.lower() == '.csv':
                    continue
                zf.write(item, item.relative_to(jury_kit_dir))
        print(f"✅ Archivio ZIP (senza CSV) creato: '{zip_path.name}'")
    except Exception as e:
        print(f"❌ ERRORE durante la creazione dello ZIP: {e}")


# --- LOGICA DI GENERAZIONE (PPT E CLASSIFICA) ---
def _build_ppt(ppt_path: Path, title_path_pairs: list[tuple[str, Path]]) -> None:
    if Presentation is None: return
    try:
        prs = Presentation()
        slide_w, slide_h = prs.slide_width, prs.slide_height
        blank_layout = prs.slide_layouts[6]
        bottom_band_h, top_margin, side_margin = Inches(0.4), Inches(0.1), Inches(0.1)
        avail_w, avail_h = slide_w - 2 * side_margin, slide_h - bottom_band_h - top_margin

        for title, img_path in title_path_pairs:
            slide = prs.slides.add_slide(blank_layout)
            try:
                # --- INIZIO CODICE CORRETTO E ROBUSTO ---
                if Image is None:
                    raise ImportError("La libreria Pillow non è installata, impossibile continuare.")

                # 1. Apri l'immagine con Pillow per ottenere le dimensioni REALI
                img_w, img_h = Image.open(img_path).size

                # 2. Aggiungi l'immagine alla slide (inizialmente con dimensioni qualsiasi)
                pic = slide.shapes.add_picture(str(img_path), left=side_margin, top=top_margin)

                # 3. Calcola la scala usando le dimensioni ottenute da Pillow
                scale = min(avail_w / img_w, avail_h / img_h)

                # 4. Applica le nuove dimensioni e posizione alla forma sulla slide
                pic.width, pic.height = int(img_w * scale), int(img_h * scale)
                pic.left, pic.top = int((slide_w - pic.width) / 2), top_margin
                # --- FINE CODICE CORRETTO E ROBUSTO ---

                tx_box = slide.shapes.add_textbox(0, slide_h - bottom_band_h, slide_w, bottom_band_h)
                tf = tx_box.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = title
                run.font.size, run.font.bold = Pt(22), True
            except Exception as e:
                print(f"ERRORE: Impossibile aggiungere l'immagine '{img_path.name}': {e}")
        prs.save(str(ppt_path))
    except Exception as e:
        print(f"\n❌ ERRORE CRITICO durante la creazione della presentazione '{ppt_path.name}': {e}")


def _build_leaderboard_ppt(ppt_path: Path, ranked_entries: list[tuple[int, str, Path]]) -> bool:
    if Presentation is None: return False
    try:
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide_w, slide_h = prs.slide_width, prs.slide_height
        for rank, title, img_path in ranked_entries:
            # Slide con testo
            title_slide = prs.slides.add_slide(blank_layout)
            tx_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), slide_w - Inches(1), slide_h - Inches(1))
            tf = tx_box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            author, _ = _parse_filename(img_path)

            p_pos = tf.paragraphs[0];
            p_pos.alignment = PP_ALIGN.CENTER
            run_pos = p_pos.add_run();
            run_pos.text = f"{rank}° Classificato"
            run_pos.font.size, run_pos.font.bold = Pt(54), True

            p_title = tf.add_paragraph();
            p_title.alignment = PP_ALIGN.CENTER
            run_title = p_title.add_run();
            run_title.text = title
            run_title.font.size = Pt(32)

            p_author = tf.add_paragraph();
            p_author.alignment = PP_ALIGN.CENTER
            run_author = p_author.add_run();
            run_author.text = f"di {author}"
            run_author.font.size, run_author.font.italic = Pt(24), True

            # Slide con immagine
            img_slide = prs.slides.add_slide(blank_layout)
            try:
                # --- INIZIO CODICE CORRETTO E ROBUSTO ---
                if Image is None:
                    raise ImportError("La libreria Pillow non è installata, impossibile continuare.")

                # 1. Apri l'immagine con Pillow per ottenere le dimensioni REALI
                img_w, img_h = Image.open(img_path).size

                # 2. Aggiungi l'immagine alla slide
                pic = img_slide.shapes.add_picture(str(img_path), 0, 0)

                # 3. Calcola la scala usando le dimensioni ottenute da Pillow
                scale = min(slide_w / img_w, slide_h / img_h)

                # 4. Applica le nuove dimensioni e la posizione
                pic.width, pic.height = int(img_w * scale), int(img_h * scale)
                pic.left, pic.top = int((slide_w - pic.width) / 2), int((slide_h - pic.height) / 2)
                # --- FINE CODICE CORRETTO E ROBUSTO ---
            except Exception as e:
                print(f"ERRORE: Impossibile aggiungere l'immagine '{img_path.name}': {e}")
        prs.save(str(ppt_path))
        return True
    except Exception as e:
        print(f"\n❌ ERRORE CRITICO durante la creazione della presentazione PowerPoint: {e}")
        return False


def _find_judge_csvs(judges_dir: Path) -> list[Path]:
    if not judges_dir.exists(): return []
    return [f for d in judges_dir.iterdir() if d.is_dir() for f in d.iterdir() if
            f.is_file() and f.suffix.lower() == ".csv"]


def _generate_leaderboard_logic(base: Path, master_csv: Path) -> None:
    judges_dir, leaderboard_dir = base / "judges", base / "leaderboard"
    leaderboard_dir.mkdir(exist_ok=True)
    all_juror_dirs = [d for d in judges_dir.iterdir() if d.is_dir()]
    judge_csvs = _find_judge_csvs(judges_dir)
    submitted_jurors = {p.parent.name for p in judge_csvs}
    missing_jurors = [d.name for d in all_juror_dirs if d.name not in submitted_jurors]

    if missing_jurors:
        print(f"⚠️  ATTENZIONE: Mancano i voti di: {', '.join(missing_jurors)}")
        if not _confirm_or_exit("Generare classifica con dati parziali?", default=False):
            print("Generazione annullata.");
            return

    criteria, titles = _read_existing_criteria(master_csv), _read_existing_titles(master_csv)
    values: dict = defaultdict(list)
    for jcsv in judge_csvs:
        try:
            with jcsv.open("r", encoding="utf-8-sig", newline="") as f:
                reader = csv.reader(f, delimiter=';')
                header = next(reader, [])
                crit_to_idx = {name: idx for idx, name in enumerate(header[1:], 1) if name}
                for row in reader:
                    if not row: continue
                    title = row[0]
                    if title in titles:
                        for crit in criteria:
                            col = crit_to_idx.get(crit)
                            if col and col < len(row) and row[col].strip():
                                try:
                                    values[(title, crit)].append(float(row[col].strip().replace(",", ".")))
                                except ValueError:
                                    print(f"AVVISO: Voto non valido in {jcsv.parent.name} per '{title}'")
        except Exception as e:
            print(f"ERRORE: Impossibile leggere '{jcsv}': {e}")

    scored = []
    for title in titles:
        means = [sum(v) / len(v) for c in criteria if (v := values.get((title, c)))]
        score = sum(means) / len(means) if means else -1.0
        cells = [f"{(sum(v) / len(v)):.2f}" if (v := values.get((title, c))) else "" for c in criteria]
        scored.append((title, cells, score))

    scored.sort(key=lambda x: (-x[2], x[0]))
    out_rows = [["Posizione", "Titolo"] + criteria + ["Media Totale"]]
    rank, last_score = 0, float('inf')
    for title, cells, total in scored:
        if total < last_score:
            rank += 1
            last_score = total
        rank_str = str(rank) if total >= 0 else "N/D"
        out_rows.append([rank_str, title] + cells + [f"{total:.2f}" if total >= 0 else "N/D"])

    leaderboard_csv = leaderboard_dir / "classifica.csv"
    try:
        with leaderboard_csv.open("w", encoding="utf-8-sig", newline="") as f:
            csv.writer(f, delimiter=';').writerows(out_rows)
        print(f"✅ Classifica CSV salvata in: {leaderboard_csv.name}")
    except (IOError, PermissionError) as e:
        print(f"❌ ERRORE: Impossibile salvare la classifica CSV: {e}");
        return

    all_ranks = sorted(list(set(int(row[0]) for row in out_rows[1:] if row[0].isdigit())))
    if not all_ranks:
        print("ℹ️ Nessun classificato valido. Salto creazione presentazione.");
        return

    rank_cutoff = all_ranks[min(9, len(all_ranks) - 1)]
    ppt_entries = []
    title_path_map = {_parse_filename(p)[1]: p for p in (base / "pictures").iterdir() if p.is_file()}
    for row in out_rows[1:]:
        if row[0].isdigit() and (current_rank := int(row[0])) <= rank_cutoff:
            if (s_title := row[1]) in title_path_map:
                ppt_entries.append((current_rank, s_title, title_path_map[s_title]))
            else:
                print(f"⚠️  AVVISO: Foto '{s_title}' non trovata per la presentazione.")

    if ppt_entries:
        ppt_entries.reverse()
        ppt_path = leaderboard_dir / "classifica.pptx"
        if _build_leaderboard_ppt(ppt_path, ppt_entries):
            print(f"✅ Presentazione classifica salvata in: {ppt_path.name}")
    else:
        print("ℹ️ Nessuna foto trovata per i primi classificati. Salto presentazione.")


# --- FUNZIONE MAIN ---

def main() -> None:
    print("---" * 15)
    print("   Benvenuto nello strumento di gestione concorsi fotografici!   ")
    print("---" * 15)

    if getattr(sys, 'frozen', False):
        base_path = Path(sys.executable).parent
    else:
        base_path = Path(__file__).resolve().parent
    root_dir = base_path / "Concorsi Orizzonti Fotografici"
    root_dir.mkdir(exist_ok=True)

    while True:
        contest_selection = _select_or_create_contest(root_dir)
        if not contest_selection:
            break

        base, folder_name = contest_selection

        if not base.exists():
            base.mkdir(parents=True)
            print(f"\n✅ Creata nuova cartella per il concorso: '{base.name}'")
        for sub in ["pictures", "jury_kit", "leaderboard", "judges"]:
            (base / sub).mkdir(exist_ok=True)

        while True:
            _display_contest_status(base, folder_name)

            choices = {
                "criteria": "📝  Gestisci Criteri di Valutazione",
                "jurors": "🛠️  Gestisci Giurati",
                "generate_kit": "🔄  Genera/Aggiorna Jury Kit Completo",
                "leaderboard": "🏆  Genera Classifica Finale",
                "open": "📂  Apri la cartella del concorso",
                "deadline": "✏️  Modifica data di scadenza",
                "switch_contest": "📚  Gestisci Concorsi (torna alla lista)",
                "help": "❓ Guida e Documentazione (online)",
                "exit": "🚪  Esci"
            }

            chosen_key = ""
            if inquirer:
                q = [inquirer.List('action', message="Menu Principale", choices=list(choices.values()), carousel=True)]
                answers = inquirer.prompt(q)
                if not answers:
                    chosen_key = "exit"
                else:
                    chosen_value = answers['action']
                    chosen_key = next(key for key, value in choices.items() if value == chosen_value)
            else:
                choice_map = list(choices.keys())
                for i, key in enumerate(choice_map, 1): print(f"  [{i}] {choices[key]}")
                while True:
                    try:
                        num = int(_input_or_exit("Scegli un'azione: "))
                        if 1 <= num <= len(choice_map):
                            chosen_key = choice_map[num - 1]
                            break
                        else:
                            print("Scelta non valida.")
                    except ValueError:
                        print("Inserisci un numero.")

            if chosen_key == "criteria":
                _action_manage_criteria(base, folder_name)
            elif chosen_key == "jurors":
                _action_manage_jurors(base / "judges")
            elif chosen_key == "generate_kit":
                _action_generate_kit(base, folder_name)
            elif chosen_key == "leaderboard":
                _action_generate_leaderboard(base, folder_name)
            elif chosen_key == "open":
                _open_folder(base)
            elif chosen_key == "deadline":
                _action_edit_deadline(root_dir, folder_name)
            elif chosen_key == "help":
                print("\n🌐 Apertura della documentazione online nel browser...")
                try:
                    webbrowser.open("https://github.com/FraMazu/orizzonti_fotografici/blob/master/README.md")
                except Exception as e:
                    print(f"❌ Impossibile aprire il browser. Dettagli: {e}")
                _input_or_exit("Premi Invio per continuare...")
            elif chosen_key == "switch_contest":
                print("\n📚 Torno alla selezione del concorso...")
                break
            elif chosen_key == "exit":
                print("\nGrazie per aver usato lo strumento. Arrivederci! 👋")
                sys.exit(0)

    print("\nGrazie per aver usato lo strumento. Arrivederci! 👋")


if __name__ == "__main__":
    if Presentation is None:
        print("AVVISO: `python-pptx` non installato, le funzioni di presentazione sono disabilitate.")
    if inquirer is None:
        print("AVVISO: `inquirer` non installato, verrà usata un'interfaccia a menu numerico.")
    if Image is None:
        print("AVVISO: `Pillow` non installato, le funzioni di generazione PPTX non funzioneranno.")

    try:
        main()
    except Exception as e:
        print("\n" + "---" * 15)
        print("❌ ERRORE IMPREVISTO: Si è verificato un problema non gestito.")
        print(f"   Dettagli tecnici: {type(e).__name__} - {e}")
        _input_or_exit("\nPremi Invio per chiudere.")
        sys.exit(1)